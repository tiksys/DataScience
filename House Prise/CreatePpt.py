# -*- coding: utf-8 -*-
"""
CreateFigures.Rで出力した図とDFを用いてpptを作成する

"""



import os
import datetime

import pandas as pd
import pptx
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.line import LineFormat


#----------関数----------
#"空のスライドを挿入する関数"
def duplicate_slide(presentation):
    title_slide_layout = presentation.slide_layouts[6]  #6は空のページ
    slide = presentation.slides.add_slide(title_slide_layout)
    return slide


#サブタイトルスライド作成関数
def create_subtitle_slide(presentation, index):
    slide = duplicate_slide(presentation)
    #text boxの位置
    text_left = Cm(1.29)
    text_top = Cm(4.52)

    #text boxの大きさ
    text_width = Cm(22.74)
    text_height = Cm(2.47)

    # 文字のフォントサイズ
    font_size = 30

    # 文字色指定(R, G, B)
    orange = RGBColor(243, 117, 0)

    #text boxの挿入
    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)

    #TOI名が入力されたとき
    if type(index) == str:
        text_box.text = "「" + index + "」について"
        par = text_box.text_frame.paragraphs[0]
        par.font.name = "ＭＳ Ｐゴシック 見出し"
    #3 or 4が入力されたとき（目次レベル）
    else:
        text_box.text = str(index) + "."
        par = text_box.text_frame.paragraphs[0]
        par.font.name = "Arial 見出し"
        run = par.add_run()
        if index == 3:
            run.text = " 全体の傾向"
        elif index == 4:
            run.text = " 調査結果"
        run.font.name = "ＭＳ Ｐゴシック 見出し"
    
    par.font.size = Pt(font_size)
    par.font.color.rgb = orange
    par.font.bold = True


#全体の傾向スライド作成関数
def create_zentaikeiko(presentation, toi_top3):
    slide = duplicate_slide(presentation)
    shapes = slide.shapes

    # 文字のフォント
    title_font_size = Pt(24)
    text_font_size = Pt(14)
    font_name = "Meiryo"

    # 文字色指定(R, G, B)
    orange = RGBColor(243, 117, 0)

    toi1 = toi_top3["TOI.Name"][0]
    toi2 = toi_top3["TOI.Name"][1]
    toi3 = toi_top3["TOI.Name"][2]

    aoi1 = toi_top3["Top.AOI"][0]
    aoi2 = toi_top3["Top.AOI"][1]
    aoi3 = toi_top3["Top.AOI"][2]

    title = "全体の傾向"
    text = "タスクの所要時間で最も差があるのは「{}」、次いで「{}」「{}」であった。".format(toi1, toi2, toi3)\
    +"\n上記のTOIで最も見た時間に差があった箇所は「{}」においては「{}」、「{}」においては「{}」、「{}」においては「{}」であった。".format(toi1, aoi1, toi2, aoi2, toi3, aoi3)

    #text box(スライドタイトル)の挿入(left,top,width,height)
    text_box = shapes.add_textbox(Cm(0.7), Cm(1.02), Cm(8.25), Cm(1.25))
    text_box.text = title
    par = text_box.text_frame.paragraphs[0]
    par.font.name = font_name
    par.font.size = title_font_size
    par.font.color.rgb = orange
    par.font.bold = True

    #タイトル下の線(left,top,width,height)
    shape = shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Cm(0.7), Cm(2.28), Cm(5.9), Cm(0))
    line = LineFormat(shape)
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Pt(1.5)
    shape.shadow.inherit = False

    #text box(スライドタイトル)の挿入(left,top,width,height)
    text_box = shapes.add_textbox(Cm(0.7), Cm(2.42), Cm(24), Cm(2.05))
    text_box.text_frame.word_wrap = True
    text_box.text = text
    par1 = text_box.text_frame.paragraphs[0]
    par2 = text_box.text_frame.paragraphs[1]
    par1.font.name = font_name
    par1.font.size = text_font_size
    par2.font.name = font_name
    par2.font.size = text_font_size
    
    return(slide)


#調査結果スライド作成関数
def create_chosa(presentation, df, toi_num, index):
    slide = duplicate_slide(presentation)
    shapes = slide.shapes

    # 文字のフォント
    title_font_size = Pt(24)
    text_font_size = Pt(14)
    font_name = "Meiryo"

    # 文字色指定(R, G, B)
    orange = RGBColor(243, 117, 0)
    
    target_row = df.iloc[toi_num, :]

    #対象TOI名
    toi = target_row[0]

    #内容記述用データ
    exp_time = target_row[1]
    nov_time = target_row[2]
    diff = target_row[3]
    
    #内容記述用AOI名
    posi_aoi1 = target_row[4]
    posi_aoi2 = target_row[5]
    posi_aoi3 = target_row[6]
    nega_aoi1 = target_row[7]
    nega_aoi2 = target_row[8]
    nega_aoi3 = target_row[9]

    title_list = ["所要時間の差", "見た時間のヒートマップ", "見た時間の差(熟＞非熟)", "時間差の原因(熟＞非熟)", "見た時間の差(熟＜非熟)", "時間差の原因(熟＜非熟)"]
    text = "タスクの所要時間について、熟練者は{}秒、非熟練者は{}秒で{}秒の差があった。\n両名の見た時間の差について、「時間をかけるべき箇所」のTOP3は「{}」「{}」「{}」。\nこれらの項目は、熟練者の方が●●※だから、熟練者の方が時間が長かったと考えられる。\n（※1回あたりの時間が長い、見た回数が多い、または両方（あてはまらない場合は外れ値の影響？））\n\n「時間をかけなくてよい箇所」項目のTOP3は「{}」「{}」「{}」であった。\nこれらの項目は、熟練者の方が●●※だから、熟練者の方が時間が短かったと考えられる。\n（※1回あたりの時間が短い、見た回数が少ない、または両方（あてはまらない場合は外れ値の影響？）".format(exp_time, nov_time, diff, posi_aoi1, posi_aoi2, posi_aoi3, nega_aoi1, nega_aoi2, nega_aoi3)

    #text box(スライドタイトル)の挿入(left,top,width,height)
    text_box = shapes.add_textbox(Cm(0.7), Cm(1.02), Cm(8.25), Cm(1.25))
    text_box.text = "「" + toi + "」について：" + title_list[index]
    par = text_box.text_frame.paragraphs[0]
    par.font.name = font_name
    par.font.size = title_font_size
    par.font.color.rgb = orange
    par.font.bold = True

    #タイトル下の線(left,top,width,height)
    shape = shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Cm(0.7), Cm(2.28), Cm(5.9), Cm(0))
    line = LineFormat(shape)
    line.color.rgb = RGBColor(0, 0, 0)
    line.width = Pt(1.5)
    shape.shadow.inherit = False

    #text box(スライドタイトル)の挿入(left,top,width,height)
    text_box = shapes.add_textbox(Cm(0.7), Cm(2.42), Cm(24), Cm(5.64))
    text_box.text_frame.word_wrap = True
    text_box.text = text
    for i in range(len(text_box.text_frame.paragraphs)):
        par = text_box.text_frame.paragraphs[i]
        par.font.name = font_name
        par.font.size = text_font_size
        if i == 0:
            par.font.color.rgb = RGBColor(0,176,80)
    return(slide)


#------------main------------
# templateとなるpptxファイルを指定する。
template_path = "CPPPackage_template.pptx"

#内容記述に用いるDFの読み込み
df = pd.read_csv("for_ppt.csv")
df = df.drop("Unnamed: 0", axis=1)

#記述用にDFを編集
toi_top3 = df
toi_top3["Abs_Diff.TOI"] = abs(toi_top3["Diff.TOI"]) 
toi_top3 = df.sort_values(by="Abs_Diff.TOI", ascending=False).reset_index(drop=True)
delete = ["Exp.TOI", "Nov.TOI", "Posi2.AOI", "Posi3.AOI", "Nega2.AOI", "Nega3.AOI", "Top.AOI", "Abs_Diff.TOI"]
toi_table = df.drop(delete, axis=1)

#templateファイルを読み込み
presentation = pptx.Presentation(template_path)

# 文字のフォント
text_font_size = Pt(14)
table_font_size = Pt(16)
font_name = "Meiryo"

#----------3. 全体の傾向サブタイトルスライドを挿入----------
create_subtitle_slide(presentation, 3)




# ----------全体の傾向1枚目----------
#スライドを作成
slide = create_zentaikeiko(presentation, toi_top3)

#画像を挿入(filepath,left,top,width,height)
slide.shapes.add_picture("output/task_duration.png", Cm(0.7), Cm(5.1), Cm(24), Cm(6.93))

#tektboxの挿入(left,top,width,height)
text_box = slide.shapes.add_textbox(Cm(0.7), Cm(12.51), Cm(24), Cm(5.56))
text_list = [
    "見た時間に差がある箇所は、以下のように解釈できる。", 
    ["\n熟＞非熟の場合、「", "熟練者のように時間をかけるべきだが、非熟練者は時間をかけていない", "（以下、時間をかけるべき箇所）」"],
    ["\n熟＜非熟の場合、「", "非熟練者は必要以上に時間をかけすぎている", "（以下、時間をかけなくてよい箇所） 」"],
    "\nこれらは、単純に見た時間の差が作業時間に与える影響があるだけでなく、非熟練者の理解や作業プロセスに問題があることを示唆している可能性がある。",
    "\n",
    "\n次頁以降は、各作業内容ごとに「時間をかけるべき箇所」「時間をかけなくてよい箇所」を洗い出していく。"
]

#文字の折り返し
text_box.text_frame.word_wrap = True

#basetext(段落の先頭のみ)の作成
base_text = ""
for i in range(0, len(text_list)):
    if type(text_list[i]) == list:
        base_text = base_text + text_list[i][0]
    else:
        base_text = base_text + text_list[i]

#textboxにbasetextを記入
text_box.text = base_text

for i in range(0, len(text_list)):
    #間で文字色の変更が必要な場合
    if type(text_list[i]) == list:
        par = text_box.text_frame.paragraphs[i]
        run = par.add_run()
        run.text = text_list[i][1]
        #文字色を緑に設定
        run.font.color.rgb = RGBColor(0,176,80)
        run = par.add_run()
        run.text = text_list[i][2]
    else:
        par = text_box.text_frame.paragraphs[i]
        
    par.font.size = text_font_size
    par.font.name = font_name

    
    
    
# ----------全体の傾向2枚目----------
#スライドを作成
slide = create_zentaikeiko(presentation, toi_top3)

#テーブルを挿入(row,col,left,top,width,height)
colmuns = [
    "TOI", "時間差(熟-非熟)",
    "時間をかけるべき箇所\n(最も差が大きい箇所)",
    "時間をかけなくてよい箇所\n(最も差が大きい箇所)"
    ]
col_width = [3, 5, 7, 7]

#tableを挿入
table_shape = slide.shapes.add_table(len(toi_table)+1, len(colmuns), Cm(1.14), Cm(5.02), Cm(23), Cm(12))

#tableオブジェクト
table = table_shape.table

#列名記入,幅設定
for i in range(len(colmuns)):
    #列幅設定
    col = table.columns[i]
    col.word_wrap = True
    col.width = Cm(col_width[i])
    
    #列名記入
    cell = table.cell(0, i)#cellオブジェクトの取得
    cell.text = colmuns[i]#textプロパティで値を設定する
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(243, 117, 0)
    par = cell.text_frame.paragraphs[0]
    par.font.size = table_font_size
    par.font.name = font_name
    if i ==  2 or i == 3:
        par = cell.text_frame.paragraphs[1]
        par.font.size = table_font_size
        par.font.name = font_name

#内容記入
for i in range(len(toi_table)):
    #列幅設定
    row = table.rows[i+1]
    row.hight = Cm(1)
    
    #行色
    if i%2 == 0:
        row_color = RGBColor(250, 214, 203)
    else:
        row_color = RGBColor(253, 236, 231)
        
    for k in range(len(colmuns)):
        cell = table.cell(i+1, k)#cellオブジェクトの取得
        cell.text = str(toi_table.iloc[i,k])#textプロパティで値を設定する
        cell.fill.solid()
        cell.fill.fore_color.rgb = row_color
        par = cell.text_frame.paragraphs[0]
        par.font.size = table_font_size
        par.font.name = font_name
        
        
        
        
#----------4. 調査結果サブタイトルスライドを挿入----------
create_subtitle_slide(presentation, 4)




#----------調査結果スライドをすべて作成----------
for i in range(len(df)):
    #対象TOI名
    toi = df["TOI.Name"][i]
    
    #サブタイトルスライド作成
    create_subtitle_slide(presentation, toi)
    
    #調査結果スライドの種類の分ループ
    for j in range(0,6):
        #調査結果スライド作成
        slide = create_chosa(presentation, df, i, j)
        
        #種類に応じて画像挿入
        if j == 0:
            #画像を挿入(filepath,left,top,width,height)
            filepath = "output/" + toi + "_01_task_duration.png"
            slide.shapes.add_picture(filepath, Cm(0.7), Cm(9), Cm(20.79), Cm(6.55))
        
        elif j == 1:
            #画像タイトル
            text_box = slide.shapes.add_textbox(Cm(4.7), Cm(8.7), Cm(3), Cm(0.86))
            text_box.text = "【熟練者】"
            text_box = slide.shapes.add_textbox(Cm(16.26), Cm(8.65), Cm(3.64), Cm(0.86))
            text_box.text = "【非熟練者】"
            
            #画像を挿入(filepath,left,top,width,height)
            filepath = "output/" + toi + "_exp_heatmap.png"
            slide.shapes.add_picture(filepath, Cm(0.7), Cm(9.62), Cm(11), Cm(5.74))
            filepath = "output/" + toi + "_nov_heatmap.png"
            slide.shapes.add_picture(filepath, Cm(12.59), Cm(9.62), Cm(11), Cm(5.74))
        
        elif j == 2:
            #画像を挿入(filepath,left,top,width,height)
            filepath = "output/" + toi + "_02_posi_fd.png"
            slide.shapes.add_picture(filepath, Cm(5.37), Cm(8.05), Cm(14.67), Cm(11))
            
        elif j == 3:
            #画像を挿入(filepath,left,top,width,height)
            filepath = "output/" + toi + "_03_posi_vc.png"
            slide.shapes.add_picture(filepath, Cm(0.7), Cm(8.53), Cm(12), Cm(9))
            filepath = "output/" + toi +  "_04_posi_median_fd.png"
            slide.shapes.add_picture(filepath, Cm(12.7), Cm(8.53), Cm(12), Cm(9))
            
        elif j == 4:
            #画像を挿入(filepath,left,top,width,height)
            filepath = "output/" + toi + "_05_nega_fd.png"
            slide.shapes.add_picture(filepath, Cm(5.37), Cm(8.05), Cm(14.67), Cm(11))
            
        elif j == 5:
            #画像を挿入(filepath,left,top,width,height)
            filepath = "output/" + toi + "_06_nega_vc.png"
            slide.shapes.add_picture(filepath, Cm(0.7), Cm(8.53), Cm(12), Cm(9))
            filepath = "output/" + toi +  "_06_nega_median_fd.png"
            slide.shapes.add_picture(filepath, Cm(12.7), Cm(8.53), Cm(12), Cm(9))

#output出力用のフォルダを作る
os.makedirs('report', exist_ok=True)

#現在時刻を取得し、データを保存する
now = datetime.datetime.now().strftime('%Y-%m-%d_%H-%M-%S')
presentation.save("report/output_{}.pptx".format(now))